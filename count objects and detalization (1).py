#!/usr/bin/env python
# coding: utf-8

# In[1]:


import pandas as pd
from pptx.util import Pt
from pptx import Presentation 
from pptx.chart.data import CategoryChartData
import win32com.client as win32


# # xlsx file processing (splitting objects by periods)

# In[2]:


objects = pd.read_excel(r'C:\Users\yav\Desktop\автоматизация отчета о детализации объектов\4,1,2 кв 18.07.xlsx', header = 1)
objects['date_act'] = pd.to_datetime(objects['Дата акта'],dayfirst = True)
objects['dont_active_data'] = pd.to_datetime(objects['Дата прекращения деятельности субъекта'],dayfirst = True)
objects['date_go'] = pd.to_datetime(objects['Дата последнего обхода/АБО'],dayfirst = True)
objects['quarter'] = pd.PeriodIndex(objects['date_go'], freq = 'Q')
all_list_objects = []
objects_294 = objects[(objects['Наименование группы мониторинга'].str.contains('294', na = False))]
objects_1q = objects[(objects['quarter'] == '2022Q1')]
objects_2q = objects[(objects['quarter'] == '2022Q2')]
objects_323 = objects[(objects['Наименование группы мониторинга'].str.contains('323', na = False))]
all_list_objects.append(objects_294)
all_list_objects.append(objects_1q)
all_list_objects.append(objects_2q)
all_list_objects.append(objects_323)
all_list_objects.append(objects)


# # detalization

# In[3]:


def detalization(objects):

    vno_nc_recomendovan = objects[(objects['Статус БП последнего обхода/АБО'] == 'ДЭПиР. Рекомендовано к публикации')]
    vno_nc_vidjt = objects[(objects['Статус БП последнего обхода/АБО'] == 'ДЭПиР. Отправлено в виджет')]

    
    vno_nc_act = objects[(objects['Статус БП последнего обхода/АБО'] == 'Информация актуализирована')]
    vno_nc_act = vno_nc_act[(vno_nc_act['Статус для целей обложения ТС'].str.contains('Выяв', na = False)) 
                  |  (vno_nc_act['Статус для целей обложения ТС'].str.contains('Нед', na = False))]

    vno_nc_act_cound_days = vno_nc_act[(vno_nc_act['Статус для целей обложения ТС'] != 'Выявление нового объекта (подлежит постановке на учет в ФНС) по п.5.10.2(1) (в связи с отменой Акта)')]
    vno_nc_act_cound_days = vno_nc_act_cound_days[(vno_nc_act_cound_days['Статус для целей обложения ТС'] != 'Выявление нового объекта (подлежит постановке на учет в ФНС) по п.5.10.2 (в связи с отменой Акта)')]
    vno_nc_act_cound_days = vno_nc_act_cound_days[(vno_nc_act_cound_days['Статус для целей обложения ТС'] != 'Выявление нового объекта (подлежит постановке на учет в ФНС) по п.5.10.1 (в связи с исключением из Списка объектов)')]
    vno_nc_act_cound_days = vno_nc_act_cound_days['date_act'] - vno_nc_act_cound_days['date_go']




    ia = objects[(objects['Статус БП последнего обхода/АБО'] == 'Информация актуализирована')]
    ia = ia[ (ia['Статус для целей обложения ТС'].str.contains('Под', na = False))
           | (ia['Статус для целей обложения ТС'].str.contains('пат', na = False))
           | (ia['Статус для целей обложения ТС'].str.contains('ЕСХ', na = False))]

    yv = objects[(objects['Статус для целей обложения ТС'].str.contains('Под', na = False))]
    yv = yv[(yv['Статус БП последнего обхода/АБО'] == 'Информация актуализирована')]

    psn = objects[(objects['Статус для целей обложения ТС'].str.contains('пате', na = False))]
    psn = psn[(psn['Статус БП последнего обхода/АБО'] == 'Информация актуализирована')]

    esxn = objects[(objects['Статус для целей обложения ТС'].str.contains('ЕСХ', na = False))]
    esxn = esxn[(esxn['Статус БП последнего обхода/АБО'] == 'Информация актуализирована')]


    cnd = objects[(objects['Статус БП последнего обхода/АБО'].str.contains('ЦНД.', na = False))]


    mkmcn_request = objects[(objects['Статус БП последнего обхода/АБО'].str.contains('МКМЦН.', na = False))
                   & (objects['Наименование ТО'].str.contains('[(]', na = False))]
    mkmcn = objects[(objects['Статус БП последнего обхода/АБО'].str.contains('МКМЦН.', na = False))]


    refusal_all = objects[(objects['Статус БП последнего обхода/АБО'].str.contains('Информа', na = False))
                & (objects['Статус для целей обложения ТС'].str.contains('нет ХС', na = False))
                & ((objects['Статус результата обхода'].str.contains('Отказ', na = False))
                | (objects['Статус результата обхода'].str.contains('ХС не установлен', na = False))     
                | (objects['Статус результата обхода'].str.contains('Не подтверж', na = False)))]

    refusal = objects[(objects['Статус результата обхода'].str.contains('Отказ', na = False))]
    refusal = refusal[(refusal['Статус БП последнего обхода/АБО'] == 'Информация актуализирована')]
    refusal = refusal[(refusal['Статус для целей обложения ТС'] == 'Обложение ТС - нет ХС')]

    close = objects[(objects['Статус результата обхода'].str.contains('Закры', na = False))]
    close = close[(close['Статус для целей обложения ТС'] == 'Обложение ТС - нет ХС')]
    close = close[(close['Статус БП последнего обхода/АБО'] == 'Информация актуализирована')]

    dont_active = objects[(objects['Статус результата обхода'].str.contains('Недейст', na = False))]
    dont_active = dont_active[(dont_active['Статус БП последнего обхода/АБО'] == 'Информация актуализирована')]
    dont_active = dont_active[(dont_active['Статус для целей обложения ТС'] == 'Обложение ТС - нет ХС')]

    dont_active_befor_go = dont_active[dont_active['dont_active_data'] < dont_active['date_go']]

    dont_confirmed_subject = objects[(objects['Статус результата обхода'].str.contains('Не подтверж', na = False))]
    dont_confirmed_subject = dont_confirmed_subject[(dont_confirmed_subject['Статус БП последнего обхода/АБО'] == 'Информация актуализирована')]
    dont_confirmed_subject = dont_confirmed_subject[(dont_confirmed_subject['Статус для целей обложения ТС'] == 'Обложение ТС - ошибочный ХС')]

    subject_not_full_identivication  = objects[(objects['Статус результата обхода'].str.contains('ХС не установлен', na = False))]
    subject_not_full_identivication = subject_not_full_identivication[(subject_not_full_identivication['Статус БП последнего обхода/АБО'] == 'Информация актуализирована')]
    subject_not_full_identivication = subject_not_full_identivication[(subject_not_full_identivication['Статус для целей обложения ТС'] == 'Обложение ТС - нет ХС')]

    cancel = objects[(objects['Статус БП последнего обхода/АБО'].str.contains('Отме', na = False))]

    trade_completed = objects[(objects['Статус БП последнего обхода/АБО'].str.contains('Информа', na = False))]
    trade_completed  = trade_completed[(trade_completed['Статус результата обхода'].str.contains('Прекращ', na = False))]
    trade_completed = trade_completed[(trade_completed['Статус для целей обложения ТС'] == 'Обложение ТС - нет объекта обложения')]

    food_objects = objects[(objects['Статус БП последнего обхода/АБО'].str.contains('Информа', na = False))]
    food_objects  = food_objects[(food_objects['Статус результата обхода'].str.contains('общеп', na = False))]
    food_objects = food_objects[(food_objects['Статус для целей обложения ТС'] == 'Обложение ТС - нет объекта обложения')]

    storage = objects[(objects['Статус БП последнего обхода/АБО'].str.contains('Информа', na = False))]
    storage  = storage[(storage['Статус результата обхода'].str.contains('скла', na = False))]
    storage = storage[(storage['Статус для целей обложения ТС'] == 'Обложение ТС - нет объекта обложения')]

    pawnshop = objects[(objects['Статус БП последнего обхода/АБО'].str.contains('Информа', na = False))]
    pawnshop  = pawnshop[(pawnshop['Статус результата обхода'].str.contains('ломба', na = False))]
    pawnshop = pawnshop[(pawnshop['Статус для целей обложения ТС'] == 'Обложение ТС - нет объекта обложения')]

    return [vno_nc_recomendovan.count()[0], 
            vno_nc_vidjt.count()[0], 
            vno_nc_act.count()[0], 
            ia.count()[0],
            yv.count()[0], 
            psn.count()[0], 
            esxn.count()[0],
            cnd.count()[0],
            mkmcn.count()[0], 
            mkmcn_request.count()[0],
            refusal.count()[0],
            close.count()[0], 
            dont_active.count()[0],
            dont_confirmed_subject.count()[0],
            subject_not_full_identivication.count()[0], 
            cancel.count()[0], 
            trade_completed.count()[0], 
            food_objects.count()[0], 
            storage.count()[0], 
            pawnshop.count()[0],
            dont_active_befor_go.count()[0],
            vno_nc_act_cound_days]

detalizion_list = []
for objects in all_list_objects:
    objects_all = {}
    
    objects_all['vno_nc_recomendovan'] = detalization(objects)[0]
    objects_all['vno_nc_vidjt'] = detalization(objects)[1]
    objects_all['vno_nc_act'] = detalization(objects)[2]
    objects_all['total_vno'] = detalization(objects)[0]+ detalization(objects)[1]+ detalization(objects)[2]
    objects_all['ia'] = detalization(objects)[3]
    objects_all['yv'] = detalization(objects)[4]
    objects_all['psn'] = detalization(objects)[5]
    objects_all['esxn'] = detalization(objects)[6]
    objects_all['cnd'] = detalization(objects)[7]
    objects_all['mkmcn'] = detalization(objects)[8]
    objects_all['mkmcn_request'] = detalization(objects)[9]
    objects_all['mkmcn_work'] = detalization(objects)[8] - detalization(objects)[9]
    objects_all['refusal'] = detalization(objects)[10]
    objects_all['close'] = detalization(objects)[11]
    objects_all['dont_active'] = detalization(objects)[12]
    objects_all['dont_confirmed_subject'] = detalization(objects)[13]
    objects_all['subject_not_full_identivication'] = detalization(objects)[14]
    objects_all['refusal_all'] = detalization(objects)[10] + detalization(objects)[13] + detalization(objects)[14]
    objects_all['cancel'] = detalization(objects)[15]
    objects_all['trade_completed'] = detalization(objects)[16]
    objects_all['food_objects'] =  detalization(objects)[17]
    objects_all['storage'] = detalization(objects)[18]
    objects_all['pawnshop'] = detalization(objects)[19]
    objects_all['dont_active_befor_go'] = detalization(objects)[20]
    objects_all['dont_active_after_go'] = detalization(objects)[12] - detalization(objects)[20]
    q = detalization(objects)[21]
    
    objects_all['min_date'] = str(q.min()).split(' ')[0]
    objects_all['max_date'] = str(q.max()).split(' ')[0]
    objects_all['mean_date'] = str(q.mean()).split(' ')[0]
    
    detalizion_list.append(objects_all)
df = pd.DataFrame(detalizion_list, index=['294', '1q', '2q', '323', 'Total_all'])
df = df.T


# In[4]:


df


# In[5]:


rez = df.loc[['vno_nc_recomendovan', 'vno_nc_vidjt', 'vno_nc_act', 'ia', 'cnd', 'mkmcn', 'close',
       'dont_active', 'refusal_all']].sum()
rez2 = df.loc[['vno_nc_recomendovan', 'vno_nc_vidjt', 'vno_nc_act', 'ia', 'cnd', 'mkmcn','close',
       'dont_active', 'refusal_all', 'cancel',
       'trade_completed', 'food_objects', 'storage', 'pawnshop',
       ]].sum()
rez3 = df.loc[['total_vno','vno_nc_recomendovan', 'vno_nc_vidjt', 'vno_nc_act', 'ia', 'yv', 'psn',
       'esxn', 'cnd', 'mkmcn', 'mkmcn_request', 'mkmcn_work', 'refusal_all', 'refusal', 'close',
        'subject_not_full_identivication', 'dont_confirmed_subject', 'dont_active','dont_active_befor_go',
    'dont_active_after_go', 'cancel', 'trade_completed', 'food_objects', 'storage', 'pawnshop']]
df = df.loc[['min_date', 'max_date', 'mean_date']].replace({'NaT': '0'})
rez4 = df.loc[['min_date', 'max_date', 'mean_date']].astype(int)
plan_2q = 3750
plan_1q = 5850
plan_4q = 5042
rez3['2q%'] = rez3['2q'] / rez[2] * 100
rez3['2q%'] = rez3['2q%'].astype(float).round(1)
rez3['1q%'] = rez3['1q'] / rez[1] * 100
rez3['1q%'] = rez3['1q%'].astype(float).round(1)
rez3['294%'] = rez3['294'] / rez[0] * 100
rez3['294%'] = rez3['294%'].astype(float).round(1)
rez3['323%'] = rez3['323'] / rez[3] * 100
rez3['323%'] = rez3['323%'].astype(float).round(1)
cols = list(rez3)
cols.insert(1, cols.pop(cols.index('294%')))
cols.insert(3, cols.pop(cols.index('1q%')))
cols.insert(5, cols.pop(cols.index('2q%')))
cols.insert(7, cols.pop(cols.index('323%')))
rez3 = rez3.loc[:, cols]
rez3 = rez3.append(rez4)
rez3 = rez3.T
rez3['Total_for_slide'] = rez
rez3['Total_all'] = rez2
rez3 = rez3.T
rez3['2q%']['Total_for_slide'] = (rez3['2q']['Total_for_slide'] / plan_2q *100)
rez3['2q%']['Total_for_slide'] = round(rez3['2q%']['Total_for_slide'],1)
rez3['1q%']['Total_for_slide'] = (rez3['1q']['Total_for_slide'] / plan_1q *100)
rez3['1q%']['Total_for_slide'] = round(rez3['1q%']['Total_for_slide'],1)
rez3['294%']['Total_for_slide'] = (rez3['294']['Total_for_slide'] / plan_4q *100)
rez3['294%']['Total_for_slide'] = round(rez3['294%']['Total_for_slide'],1)
rez3['2q']['Total_for_slide'] = round(rez3['2q']['Total_for_slide'])
rez3['1q']['Total_for_slide'] = round(rez3['1q']['Total_for_slide'])


# In[6]:


rez3


# In[7]:


rez3[['1q','1q%']]


# # exprot data in presentation

# In[23]:


ppt=Presentation(r'C:\Users\yav\Desktop\автоматизация отчета о детализации объектов\04072022 Результаты рассмотрения обходов.pptx')
s0 = ppt.slides[0]
s1 = ppt.slides[1]
s2 = ppt.slides[2]
s3 = ppt.slides[3]
datenow = pd.Timestamp.now() 
datenow = datenow.strftime("%d.%m.%Y")

for shape in s0.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if 'datenow' in run.text:
                run.text = run.text.replace('datenow', datenow)
                font = run.font
                font.size = Pt(16)


for shape in s1.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if 'Total_for_slide' in run.text:
                run.text = run.text.replace('Total_for_slide', str(rez3['1q']['Total_for_slide']))
                font = run.font
                font.size = Pt(16)
            elif 'Total_per' in run.text:
                run.text = run.text.replace('Total_per', str(rez3['1q%']['Total_for_slide']))
                font = run.font
                font.size = Pt(16)
            elif '_3' in run.text:
                run.text = run.text.replace('_3', str(rez3['1q']['total_vno']))
                font = run.font
                font.size = Pt(11)
            elif '_4' in run.text:
                run.text = run.text.replace('_4', str(rez3['1q%']['total_vno']))
                font = run.font
                font.size = Pt(11)
            elif '_5' in run.text:
                run.text = run.text.replace('_5', str(rez3['1q']['mean_date']))
                font = run.font
                font.size = Pt(11)
            elif '_6' in run.text:
                run.text = run.text.replace('_6', str(rez3['1q']['min_date']))
                font = run.font
                font.size = Pt(11)
            elif '_7' in run.text:
                run.text = run.text.replace('_7', str(rez3['1q']['max_date']))
                font = run.font
                font.size = Pt(11)
            elif '_8' in run.text:
                run.text = run.text.replace('_8', f" ({rez3['1q%']['ia']}%)")
                font = run.font
                font.size = Pt(14)
            elif '_9' in run.text:
                run.text = run.text.replace('_9', f" ({rez3['1q%']['vno_nc_vidjt']}%)")
                font = run.font
                font.size = Pt(14)
            elif '1.0' in run.text:
                run.text = run.text.replace('1.0', f" ({rez3['1q%']['vno_nc_act']}%)")
                font = run.font
                font.size = Pt(14)
            elif '1.1' in run.text:
                run.text = run.text.replace('1.1', f" ({rez3['1q%']['refusal_all']}%)")
                font = run.font
                font.size = Pt(14)
            elif '1.2' in run.text:
                run.text = run.text.replace('1.2', f" ({rez3['1q%']['close']}%)")
                font = run.font
                font.size = Pt(14)
            elif '1.3' in run.text:
                run.text = run.text.replace('1.3', f" ({rez3['1q%']['dont_active_befor_go']}%)")
                font = run.font
                font.size = Pt(14)
            elif '1.4' in run.text:
                run.text = run.text.replace('1.4', f" ({rez3['1q%']['dont_active_after_go']}%)")
                font = run.font
                font.size = Pt(14)
            elif '1.5' in run.text:
                run.text = run.text.replace('1.5', f" ({rez3['1q%']['mkmcn']}%)")
                font = run.font
                font.size = Pt(14)
            elif '1.6' in run.text:
                run.text = run.text.replace('1.6', f" ({rez3['1q%']['cnd']}%)")
                font = run.font
                font.size = Pt(14)
            elif '1.7' in run.text:
                run.text = run.text.replace('1.7', f" ({rez3['1q%']['vno_nc_recomendovan']}%)")
                font = run.font
                font.size = Pt(14)
            elif 'tot' in run.text:
                run.text = run.text.replace('tot', str(plan_1q - rez3['1q']['Total_for_slide']))
                font = run.font
                font.size = Pt(16)
            elif 'per' in run.text:
                run.text = run.text.replace('per', str(round(100-rez3['1q%']['Total_for_slide'],1)))
                font = run.font
                font.size = Pt(16)
                
                
                
for shape in s2.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if '_1' in run.text:
                run.text = run.text.replace('_1', str(rez3['2q']['Total_for_slide']))
                font = run.font
                font.size = Pt(16)
            elif '_2' in run.text:
                run.text = run.text.replace('_2', str(rez3['2q%']['Total_for_slide']))
                font = run.font
                font.size = Pt(16)
            elif '_3' in run.text:
                run.text = run.text.replace('_3', str(rez3['2q']['total_vno']))
                font = run.font
                font.size = Pt(11)
            elif '_4' in run.text:
                run.text = run.text.replace('_4', f"{rez3['2q%']['total_vno']}%")
                font = run.font
                font.size = Pt(11)
            elif '5_' in run.text:
                run.text = run.text.replace('5_', f"{rez3['2q']['mkmcn']}({rez3['2q%']['mkmcn']}%)")
                font = run.font
                font.size = Pt(12)
            elif '_6' in run.text:
                run.text = run.text.replace('_6', str(rez3['2q']['mean_date']))
                font = run.font
                font.size = Pt(11)
            elif '_7' in run.text:
                run.text = run.text.replace('_7', str(rez3['2q']['min_date']))
                font = run.font
                font.size = Pt(11)
            elif '_8' in run.text:
                run.text = run.text.replace('_8', str(rez3['2q']['max_date']))
                font = run.font
                font.size = Pt(11)
            elif '(_9)' in run.text:
                run.text = run.text.replace('(_9)', f"  ({rez3['2q%']['vno_nc_recomendovan']}%)")
                font = run.font
                font.size = Pt(14)
            elif '(9.1)' in run.text:
                run.text = run.text.replace('(9.1)', f"  ({rez3['2q%']['vno_nc_vidjt']}%)")
                font = run.font
                font.size = Pt(14)
            elif '9.2)' in run.text:
                run.text = run.text.replace('9.2)', f"({rez3['2q%']['vno_nc_act']}%)")
                font = run.font
                font.size = Pt(14)
            elif '(9.3)' in run.text:
                run.text = run.text.replace('(9.3)', f" ({rez3['2q%']['refusal_all']}%)")
                font = run.font
                font.size = Pt(14)
            elif '(9.4)' in run.text:
                run.text = run.text.replace('(9.4)', f" ({rez3['2q%']['close']}%)")
                font = run.font
                font.size = Pt(14)
            elif '9.9)' in run.text:
                run.text = run.text.replace('9.9)', f" ({rez3['2q%']['dont_active_befor_go']}%)")
                font = run.font
                font.size = Pt(14)
            elif '9.5)' in run.text:
                run.text = run.text.replace('9.5)', f" ({rez3['2q%']['dont_active_after_go']}%)")
                font = run.font
                font.size = Pt(14)
            elif '9.6)' in run.text:
                run.text = run.text.replace('9.6)', f" ({rez3['2q%']['cnd']}%)")
                font = run.font
                font.size = Pt(14)
            elif '9.7)' in run.text:
                run.text = run.text.replace('9.7)', f" ({rez3['2q%']['mkmcn_work']}%)")
                font = run.font
                font.size = Pt(14)
            elif '9.8)' in run.text:
                run.text = run.text.replace('9.8)', f" ({rez3['2q%']['mkmcn_request']}%)")
                font = run.font
                font.size = Pt(14)
            elif '1.0' in run.text:
                run.text = run.text.replace('1.0', f" ({rez3['2q%']['ia']}%)")
                font = run.font
                font.size = Pt(14)
                
                
                
for shape in s3.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if 'datenow' in run.text:
                run.text = run.text.replace('datenow', datenow)
                font = run.font
                font.size = Pt(13)
            elif 'vno_nc_vidjt' in run.text:
                run.text = run.text.replace('vno_nc_vidjt', str(rez3['294']['vno_nc_vidjt']))
                font = run.font
                font.size = Pt(12)
            elif 'vno_nc_act' in run.text:
                run.text = run.text.replace('vno_nc_act', str(rez3['294']['vno_nc_act']))
                font = run.font
                font.size = Pt(12)
            elif 'vno_nc_recomendovan' in run.text:
                run.text = run.text.replace('vno_nc_recomendovan', str(rez3['294']['vno_nc_recomendovan']))
                font = run.font
                font.size = Pt(12)
            elif 'trade_completed' in run.text:
                run.text = run.text.replace('trade_completed', str(rez3['294']['trade_completed']))
                font = run.font
                font.size = Pt(12)
            elif 'close' in run.text:
                run.text = run.text.replace('close', str(rez3['294']['close']))
                font = run.font
                font.size = Pt(12)
            elif 'refusal' in run.text:
                run.text = run.text.replace('refusal', str(rez3['294']['refusal_all']))
                font = run.font
                font.size = Pt(12)
            elif 'dont_active' in run.text:
                run.text = run.text.replace('dont_active', str(rez3['294']['dont_active']))
                font = run.font
                font.size = Pt(12)
            elif 'ia' in run.text:
                run.text = run.text.replace('ia', str(rez3['294']['ia']))
                font = run.font
                font.size = Pt(12)
            elif 'work' in run.text:
                run.text = run.text.replace('work', str(rez3['294']['mkmcn']+rez3['294']['cnd']))
                font = run.font
                font.size = Pt(12)
            

chart2 = ppt.slides[1].shapes[2].chart
chart_data1 = CategoryChartData()
chart_data1.categories = ['', '', '']
chart_data1.add_series('Series 1', (rez3['1q']['mkmcn'], 
                                   rez3['1q']['cnd'], 0, 
                                   rez3['1q']['dont_active_after_go'],
                                   rez3['1q']['dont_active_befor_go'], 0, 
                                   rez3['1q']['close'], 
                                   rez3['1q']['refusal_all'], 0, 
                                   rez3['1q']['vno_nc_act'], 
                                   rez3['1q']['vno_nc_vidjt'], 
                                   rez3['1q']['vno_nc_recomendovan'], 
                                   rez3['1q']['ia']))                
                
chart3 = ppt.slides[2].shapes[2].chart
chart_data = CategoryChartData()
chart_data.categories = ['', '', '']
chart_data.add_series('Series 1', (rez3['2q']['mkmcn_request'], 
                                   rez3['2q']['mkmcn_work'], 
                                   rez3['2q']['cnd'], 0, 
                                   rez3['2q']['dont_active_after_go'],
                                   rez3['2q']['dont_active_befor_go'], 0, 
                                   rez3['2q']['close'], 
                                   rez3['2q']['refusal_all'], 0, 
                                   rez3['2q']['vno_nc_act'], 
                                   rez3['2q']['vno_nc_vidjt'], 
                                   rez3['2q']['vno_nc_recomendovan'], 0, 
                                   rez3['2q']['ia']))
chart3.replace_data(chart_data)
chart2.replace_data(chart_data1)
name_result = 'C:\\Users\\yav\\Desktop\\автоматизация отчета о детализации объектов\\'
name_result = name_result + datenow + ' ' + 'Результаты рассмотрения обходов.pptx'
ppt.save(name_result)


# # send result on email in MS Outlook

# In[33]:


outlook = win32.Dispatch('outlook.application')
mail = outlook.CreateItem(0)
mail.To = 'MedvedevRA@tax.mos.ru'
mail.Subject = datenow + ' ' + 'Результаты рассмотрения обходов.pptx'
#mail.Body = 'Message body'
mail.HTMLBody = '<h2>Добрый день, направляю результаты рассмотрения обходов по состоянию на %s</h2>' %datenow #this field is optional


attachment  = name_result
mail.Attachments.Add(attachment)

mail.Send()


# In[ ]:




