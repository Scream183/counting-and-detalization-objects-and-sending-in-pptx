#!/usr/bin/env python
# coding: utf-8

# In[1]:


import pandas as pd
from pptx.util import Pt
from pptx import Presentation 
from pptx.chart.data import CategoryChartData


# # xlsx file processing (splitting objects by periods)

# In[2]:


objects = pd.read_excel(r'C:\Users\yav\Desktop\автоматизация отчета о детализации объектов\4,1,2 кв 08.07 (2).xlsx', header = 1, dtype = str)
objects['date_act'] = pd.to_datetime(objects['Дата акта'],dayfirst = True)
objects['dont_active_data'] = pd.to_datetime(objects['Дата прекращения деятельности субъекта'],dayfirst = True)
objects['date_go'] = pd.to_datetime(objects['Дата последнего обхода/АБО'],dayfirst = True)
objects['quarter'] = pd.PeriodIndex(objects['date_go'], freq = 'Q')
all_list_objects = []
objects_4q = objects[(objects['quarter'] == '2021Q4')]
objects_1q = objects[(objects['quarter'] == '2022Q1')]
objects_2q = objects[(objects['quarter'] == '2022Q2')]
objects_294 = objects[(objects['Наименование группы мониторинга'].str.contains('294', na = False))]
objects_323 = objects[(objects['Наименование группы мониторинга'].str.contains('323', na = False))]
all_list_objects.append(objects_4q)
all_list_objects.append(objects_1q)
all_list_objects.append(objects_2q)
all_list_objects.append(objects_323)
all_list_objects.append(objects)


# # detalization

# In[3]:


def detalization(objects):

    vno_nc_recomendovan = objects[(objects['Статус БП последнего обхода/АБО'] == 'ДЭПиР. Рекомендовано к публикации')]
    vno_nc_vidjt = objects[(objects['Статус БП последнего обхода/АБО'] == 'ДЭПиР. Отправлено в виджет')]

    objects['date_act'] = pd.to_datetime(objects['Дата акта'],dayfirst = True)
    objects['dont_active_data'] = pd.to_datetime(objects['Дата прекращения деятельности субъекта'],dayfirst = True)
    objects['date_go'] = pd.to_datetime(objects['Дата последнего обхода/АБО'],dayfirst = True)

    vno_nc_act = objects[(objects['Статус БП последнего обхода/АБО'] == 'Информация актуализирована')]
    vno_nc_act = vno_nc_act[(vno_nc_act['Статус для целей обложения ТС'].str.contains('Выяв', na = False)) 
                  |  (vno_nc_act['Статус для целей обложения ТС'].str.contains('Нед', na = False))]

    vno_nc_act_cound_days = vno_nc_act[(vno_nc_act['Статус для целей обложения ТС'] != 'Выявление нового объекта (подлежит постановке на учет в ФНС) по п.5.10.2(1) (в связи с отменой Акта)')]
    vno_nc_act_cound_days = vno_nc_act_cound_days[(vno_nc_act_cound_days['Статус для целей обложения ТС'] != 'Выявление нового объекта (подлежит постановке на учет в ФНС) по п.5.10.2 (в связи с отменой Акта)')]
    vno_nc_act_cound_days = vno_nc_act_cound_days[(vno_nc_act_cound_days['Статус для целей обложения ТС'] != 'Выявление нового объекта (подлежит постановке на учет в ФНС) по п.5.10.1 (в связи с исключением из Списка объектов)')]
    vno_nc_act_cound_days = vno_nc_act_cound_days['date_act'] - vno_nc_act_cound_days['date_go']




    ia = objects[(objects['Статус БП последнего обхода/АБО'] == 'Информация актуализирована')]
    ia = ia[ (ia['Статус для целей обложения ТС'].str.contains('Под', na = False))
           | (ia['Статус для целей обложения ТС'].str.contains('пат', na = False))
           | (ia['Статус для целей обложения ТС'].str.contains('ЕСХ', na = False))]

    yv = objects[(objects['Статус для целей обложения ТС'].str.contains('Под', na = False))]
    yv = yv[(yv['Статус БП последнего обхода/АБО'] == 'Информация актуализирована')]

    psn = objects[(objects['Статус для целей обложения ТС'].str.contains('пате', na = False))]
    psn = psn[(psn['Статус БП последнего обхода/АБО'] == 'Информация актуализирована')]

    esxn = objects[(objects['Статус для целей обложения ТС'].str.contains('ЕСХ', na = False))]
    esxn = esxn[(esxn['Статус БП последнего обхода/АБО'] == 'Информация актуализирована')]


    cnd = objects[(objects['Статус БП последнего обхода/АБО'].str.contains('ЦНД.', na = False))]


    mkmcn_request = objects[(objects['Статус БП последнего обхода/АБО'].str.contains('МКМЦН.', na = False))
                   & (objects['Наименование ТО'].str.contains('[(]', na = False))]
    mkmcn = objects[(objects['Статус БП последнего обхода/АБО'].str.contains('МКМЦН.', na = False))]


    refusal_all = objects[(objects['Статус БП последнего обхода/АБО'].str.contains('Информа', na = False))
                & (objects['Статус для целей обложения ТС'].str.contains('нет ХС', na = False))
                & ((objects['Статус результата обхода'].str.contains('Отказ', na = False))
                | (objects['Статус результата обхода'].str.contains('ХС не установлен', na = False))     
                | (objects['Статус результата обхода'].str.contains('Не подтверж', na = False)))]

    refusal = objects[(objects['Статус результата обхода'].str.contains('Отказ', na = False))]
    refusal = refusal[(refusal['Статус БП последнего обхода/АБО'] == 'Информация актуализирована')]
    refusal = refusal[(refusal['Статус для целей обложения ТС'] == 'Обложение ТС - нет ХС')]

    close = objects[(objects['Статус результата обхода'].str.contains('Закры', na = False))]
    close = close[(close['Статус для целей обложения ТС'] == 'Обложение ТС - нет ХС')]
    close = close[(close['Статус БП последнего обхода/АБО'] == 'Информация актуализирована')]

    dont_active = objects[(objects['Статус результата обхода'].str.contains('Недейст', na = False))]
    dont_active = dont_active[(dont_active['Статус БП последнего обхода/АБО'] == 'Информация актуализирована')]
    dont_active = dont_active[(dont_active['Статус для целей обложения ТС'] == 'Обложение ТС - нет ХС')]

    dont_active_befor_go = dont_active[dont_active['dont_active_data'] < dont_active['date_go']]

    dont_confirmed_subject = objects[(objects['Статус результата обхода'].str.contains('Не подтверж', na = False))]
    dont_confirmed_subject = dont_confirmed_subject[(dont_confirmed_subject['Статус БП последнего обхода/АБО'] == 'Информация актуализирована')]
    dont_confirmed_subject = dont_confirmed_subject[(dont_confirmed_subject['Статус для целей обложения ТС'] == 'Обложение ТС - ошибочный ХС')]

    subject_not_full_identivication  = objects[(objects['Статус результата обхода'].str.contains('ХС не установлен', na = False))]
    subject_not_full_identivication = subject_not_full_identivication[(subject_not_full_identivication['Статус БП последнего обхода/АБО'] == 'Информация актуализирована')]
    subject_not_full_identivication = subject_not_full_identivication[(subject_not_full_identivication['Статус для целей обложения ТС'] == 'Обложение ТС - нет ХС')]

    cancel = objects[(objects['Статус БП последнего обхода/АБО'].str.contains('Отме', na = False))]

    trade_completed = objects[(objects['Статус БП последнего обхода/АБО'].str.contains('Информа', na = False))]
    trade_completed  = trade_completed[(trade_completed['Статус результата обхода'].str.contains('Прекращ', na = False))]
    trade_completed = trade_completed[(trade_completed['Статус для целей обложения ТС'] == 'Обложение ТС - нет объекта обложения')]

    food_objects = objects[(objects['Статус БП последнего обхода/АБО'].str.contains('Информа', na = False))]
    food_objects  = food_objects[(food_objects['Статус результата обхода'].str.contains('общеп', na = False))]
    food_objects = food_objects[(food_objects['Статус для целей обложения ТС'] == 'Обложение ТС - нет объекта обложения')]

    storage = objects[(objects['Статус БП последнего обхода/АБО'].str.contains('Информа', na = False))]
    storage  = storage[(storage['Статус результата обхода'].str.contains('скла', na = False))]
    storage = storage[(storage['Статус для целей обложения ТС'] == 'Обложение ТС - нет объекта обложения')]

    pawnshop = objects[(objects['Статус БП последнего обхода/АБО'].str.contains('Информа', na = False))]
    pawnshop  = pawnshop[(pawnshop['Статус результата обхода'].str.contains('ломба', na = False))]
    pawnshop = pawnshop[(pawnshop['Статус для целей обложения ТС'] == 'Обложение ТС - нет объекта обложения')]

    return [vno_nc_recomendovan.count()[0], 
            vno_nc_vidjt.count()[0], 
            vno_nc_act.count()[0], 
            ia.count()[0],
            yv.count()[0], 
            psn.count()[0], 
            esxn.count()[0],
            cnd.count()[0],
            mkmcn.count()[0], 
            mkmcn_request.count()[0],
            refusal.count()[0],
            close.count()[0], 
            dont_active.count()[0],
            dont_confirmed_subject.count()[0],
            subject_not_full_identivication.count()[0], 
            refusal_all.count()[0], 
            cancel.count()[0], 
            trade_completed.count()[0], 
            food_objects.count()[0], 
            storage.count()[0], 
            pawnshop.count()[0],
            dont_active_befor_go.count()[0],
            vno_nc_act_cound_days]

detalizion_list = []
for objects in all_list_objects:
    objects_all = {}
    
    objects_all['vno_nc_recomendovan'] = detalization(objects)[0]
    objects_all['vno_nc_vidjt'] = detalization(objects)[1]
    objects_all['vno_nc_act'] = detalization(objects)[2]
    objects_all['ia'] = detalization(objects)[3]
    objects_all['yv'] = detalization(objects)[4]
    objects_all['psn'] = detalization(objects)[5]
    objects_all['esxn'] = detalization(objects)[6]
    objects_all['cnd'] = detalization(objects)[7]
    objects_all['mkmcn'] = detalization(objects)[8]
    objects_all['mkmcn_request'] = detalization(objects)[9]
    objects_all['refusal'] = detalization(objects)[10]
    objects_all['close'] = detalization(objects)[11]
    objects_all['dont_active'] = detalization(objects)[12]
    objects_all['dont_confirmed_subject'] = detalization(objects)[13]
    objects_all['subject_not_full_identivication'] = detalization(objects)[14]
    objects_all['refusal_all'] = detalization(objects)[15]
    objects_all['cancel'] = detalization(objects)[16]
    objects_all['trade_completed'] = detalization(objects)[17]
    objects_all['food_objects'] =  detalization(objects)[18]
    objects_all['storage'] = detalization(objects)[19]
    objects_all['pawnshop'] = detalization(objects)[20]
    objects_all['dont_active_befor_go'] = detalization(objects)[21]
    q = detalization(objects)[22]

    objects_all['min_data'] = str(q.min()).split(' ')[0]
    objects_all['max_data'] = str(q.max()).split(' ')[0]
    objects_all['mean_data'] = str(q.mean()).split(' ')[0]
    
    detalizion_list.append(objects_all)


# In[10]:


detalizion_list[2]['cnd']


# # counting + % + summary

# In[4]:


total_plan = 5042

for i in detalizion_list:
    
    
    total_all = (i['vno_nc_recomendovan']+i['vno_nc_vidjt']+i['vno_nc_act']+i['ia']+i['refusal']+
                 i['close']+i['dont_active']+i['dont_confirmed_subject']+
                 i['subject_not_full_identivication']+ i['cnd']+i['mkmcn'])
    

    total_all_procent = round(total_all / total_plan *100,1)
    t = f'ВСЕГО ОБЪЕКТОВ: {total_all} {total_all_procent}%'
    
    
    c_procent = round(i['cnd']/total_all*100,1)
    c = f"цнд: {i['cnd']} ({c_procent}%)"
   
    
    m_procent = round(i['mkmcn']/total_all*100,1)
    m = f"ВСЕГО мкмцн: {i['mkmcn']} ({m_procent}%)"
    

    m_procent_request = round(i['mkmcn_request']/total_all*100,1)
    m_request = f"мкмцн на запросах: {i['mkmcn_request']} ({m_procent_request}%)"
    
    
    m_work = i['mkmcn'] - i['mkmcn_request']
    m_procent_work = round(m_work/total_all*100,1)
    m_working = f'мкмцн на доработке: {m_work} ({m_procent_work}%)'
    
    
    ref_procent = round(i['refusal_all']/total_all*100,1)
    ref = f"Отказы: {i['refusal_all']} ({ref_procent}%) <--- Отказы: {i['refusal']} + ХС не установлен: {i['subject_not_full_identivication']} + ХС не подтверджен: {i['dont_confirmed_subject']}"
    
    
    clos_procent = round(i['close']/total_all*100,1)
    clos = f"закрытые: {i['close']} ({clos_procent}%)"

    ned_procent_befor = round(i['dont_active_befor_go']/total_all*100,1)
    ned_befor_go = f"Недействующие до обхода: {i['dont_active_befor_go']} ({ned_procent_befor}%)"

    ned_after_go = i['dont_active'] - i['dont_active_befor_go']
    ned_procent_after = round(ned_after_go/total_all*100,1)
    ned_after = f'Недействующие после обхода: {ned_after_go} ({ned_procent_after}%)'

    other = f"отменены: {i['cancel']}, Торговля прекращена: {i['trade_completed']}, общепит: {i['food_objects']}, склад: {i['storage']}, ломбард: {i['pawnshop']}"

    
    total_vno = i['vno_nc_recomendovan'] + i['vno_nc_vidjt'] + i['vno_nc_act']
    total_vno_nc_procent = round(total_vno / total_all *100,1)
    total_vno_nc = f'ВСЕГО ВНО,НС: {total_vno} ({total_vno_nc_procent}%)' 
    

    recomendovan_procent = round(i['vno_nc_recomendovan'] / total_all *100,1)
    recomendovan = f"ДЭПИР. Рекомендовано к публикации: {i['vno_nc_recomendovan']} ({recomendovan_procent}%)" 
    
    
    vidjt_procent = round(i['vno_nc_vidjt'] / total_all *100,1)
    vidjt = f"ДЭПиР. Отправлено в виджет: {i['vno_nc_vidjt']} ({vidjt_procent}%)"
    
    
    act_procent = round(i['vno_nc_act'] / total_all *100,1)
    act = f"Cоставлен Акт: {i['vno_nc_act']} ({act_procent}%) <-------------------- {i['mean_data']} средний срок, min - {i['min_data']}, max - {i['max_data']}"
    

    ia_procent = round(i['ia'] / total_all * 100,1)
    ia = f"ИА: {i['ia']}({ia_procent}%) <--- Уведомления: {i['yv']} + ПСН: {i['psn']} + ЕСХН: {i['esxn']}"
    
    
    print(total_all)
    print()
    print(ia)
    print()
    print(total_vno_nc)
    print(recomendovan)
    print(vidjt)
    print(act)
    print()
    print(ref)
    print(clos)
    print()
    print(ned_befor_go)
    print(ned_after)
    print()
    print(c)
    print(m)
    print(m_request)
    print(m_working)
    print()
    print(other)
    print('-'*100)


# In[12]:


detalizion_list[2]


# # exprot data in presentation

# In[16]:


from pptx import Presentation 
from pptx.chart.data import CategoryChartData
ppt=Presentation(r'C:\Users\yav\Desktop\автоматизация отчета о детализации объектов\04072022 Результаты рассмотрения обходов.pptx')
s3 = ppt.slides[3]
for shape in s3.shapes:
    if not shape.has_text_frame:
        continue
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if '_1' in run.text:
                run.text = run.text.replace('_1', str(total_all))
                font = run.font
                font.size = Pt(16)
            elif '_2' in run.text:
                run.text = run.text.replace('_2', str(total_all_procent))
                font = run.font
                font.size = Pt(16)
            elif '_3' in run.text:
                run.text = run.text.replace('_3', str(total_vno))
                font = run.font
                font.size = Pt(11)
            elif '_4' in run.text:
                run.text = run.text.replace('_4', str(total_vno_nc_procent))
                font = run.font
                font.size = Pt(11)
            elif '5_' in run.text:
                run.text = run.text.replace('5_', f"{detalizion_list[2]['mkmcn']}({m_procent}%)")
                font = run.font
                font.size = Pt(12)
            elif '_6' in run.text:
                run.text = run.text.replace('_6', str(detalizion_list[2]['mean_data']))
                font = run.font
                font.size = Pt(11)
            elif '_7' in run.text:
                run.text = run.text.replace('_7', str(detalizion_list[2]['min_data']))
                font = run.font
                font.size = Pt(11)
            elif '_8' in run.text:
                run.text = run.text.replace('_8', str(detalizion_list[2]['max_data']))
                font = run.font
                font.size = Pt(11)
            elif '(_9)' in run.text:
                run.text = run.text.replace('(_9)', f"  ({recomendovan_procent}%)")
                font = run.font
                font.size = Pt(14)
            elif '(9.1)' in run.text:
                run.text = run.text.replace('(9.1)', f"  ({vidjt_procent}%)")
                font = run.font
                font.size = Pt(14)
            elif '9.2)' in run.text:
                run.text = run.text.replace('9.2)', f"({act_procent}%)")
                font = run.font
                font.size = Pt(14)
            elif '(9.3)' in run.text:
                run.text = run.text.replace('(9.3)', f" ({ref_procent}%)")
                font = run.font
                font.size = Pt(14)
            elif '(9.4)' in run.text:
                run.text = run.text.replace('(9.4)', f" ({clos_procent}%)")
                font = run.font
                font.size = Pt(14)
            elif '9.9)' in run.text:
                run.text = run.text.replace('9.9)', f" ({ned_procent_befor}%)")
                font = run.font
                font.size = Pt(14)
            elif '9.5)' in run.text:
                run.text = run.text.replace('9.5)', f" ({ned_procent_after}%)")
                font = run.font
                font.size = Pt(14)
            elif '9.6)' in run.text:
                run.text = run.text.replace('9.6)', f" ({c_procent}%)")
                font = run.font
                font.size = Pt(14)
            elif '9.7)' in run.text:
                run.text = run.text.replace('9.7)', f" ({m_procent_work}%)")
                font = run.font
                font.size = Pt(14)
            elif '9.8)' in run.text:
                run.text = run.text.replace('9.8)', f" ({m_procent_request}%)")
                font = run.font
                font.size = Pt(14)
            elif '1.0' in run.text:
                run.text = run.text.replace('1.0', f" ({ia_procent}%)")
                font = run.font
                font.size = Pt(14)
                
                
chart = ppt.slides[3].shapes[2].chart
chart_data = CategoryChartData()
chart_data.categories = ['', '', '']
chart_data.add_series('Series 1', (detalizion_list[2]['mkmcn_request'], m_work, detalizion_list[2]['cnd'], 0, ned_after_go, detalizion_list[2]['dont_active_befor_go'], 0, detalizion_list[2]['close'], detalizion_list[2]['refusal_all'], 0, detalizion_list[2]['vno_nc_act'], detalizion_list[2]['vno_nc_vidjt'], detalizion_list[2]['vno_nc_recomendovan'], 0, detalizion_list[2]['ia']))
chart.replace_data(chart_data)
ppt.save(r'C:\Users\yav\Desktop\автоматизация отчета о детализации объектов\result.pptx')


# In[ ]:





# In[ ]:




